"""
ocr_engine_enhanced.py
──────────────────────
Text extraction engine supporting multiple document formats.

ARCHITECTURE:
  PDF / DOCX / PPTX / XLSX / EPUB / TXT
    → PyMuPDF (fitz) — primary, fast, no ML needed
    → python-docx / python-pptx / openpyxl — fallback when fitz text layer is empty

  Images (JPG / PNG / BMP / TIFF)
    → PaddleOCR — ML-based, handles handwriting and scanned images

WHY PYMUPDF FIRST:
  PaddleOCR is slow to initialise (~10 s) and unnecessary for digital PDFs
  which already have a text layer. PyMuPDF extracts that layer in milliseconds
  with 100 % accuracy, and also preserves font metadata (bold, italic, size).
"""

import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Dict, Optional

import fitz  # PyMuPDF


# ── Data model ────────────────────────────────────────────────────────────────

@dataclass
class TextBlock:
    """One span of text with its formatting metadata."""
    text: str
    confidence: float
    bbox: List                  # [x0, y0, x1, y1]
    font_size: float = 12.0
    font_name: str = "Unknown"
    is_bold: bool = False
    is_italic: bool = False
    alignment: str = "left"     # left | center | right | justify


# ── Engine ────────────────────────────────────────────────────────────────────

class EnhancedOCREngine:
    """
    Multi-format text extraction engine.
    PaddleOCR is initialised lazily — only when an image file is uploaded,
    so the heavy ML model does not slow down PDF/DOCX processing.
    """

    IMAGE_EXTS  = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"}
    FITZ_EXTS   = {".pdf", ".epub", ".xps", ".svg", ".txt", ".fb2"}
    OFFICE_EXTS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}

    def __init__(self, languages: Optional[List[str]] = None):
        self.languages = languages or ["en"]
        self._paddle_ocr = None   # lazy init

    # ── Public entry point ────────────────────────────────────────────────────

    def extract_text(self, file_path: str) -> Dict:
        """
        Detect file type and dispatch to the right extractor.
        Always returns a dict with at least:
          raw_text, blocks, confidence, block_count, source
        """
        ext = Path(file_path).suffix.lower()

        if ext in self.IMAGE_EXTS:
            return self._extract_image(file_path)

        if ext in self.FITZ_EXTS or ext in self.OFFICE_EXTS:
            result = self._extract_fitz(file_path)
            # If fitz returned nothing, try office-specific fallbacks
            if not result["raw_text"].strip():
                return self._office_fallback(file_path, ext)
            return result

        # Unknown — try fitz anyway
        return self._extract_fitz(file_path)

    # ── PyMuPDF extraction ────────────────────────────────────────────────────

    def _extract_fitz(self, file_path: str) -> Dict:
        try:
            doc = fitz.open(file_path)
        except Exception as e:
            return self._error_result(f"PyMuPDF could not open file: {e}")

        all_blocks: List[TextBlock] = []
        page_data = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_dict = page.get_text("dict")
            page_text = ""
            page_blocks: List[TextBlock] = []

            for block in page_dict.get("blocks", []):
                if "lines" not in block:
                    continue
                for line in block["lines"]:
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if not text:
                            continue

                        bbox      = list(span.get("bbox", [0, 0, 0, 0]))
                        font_size = float(span.get("size", 12))
                        font_name = span.get("font", "Unknown")
                        flags     = span.get("flags", 0)
                        is_bold   = bool(flags & 0b0001)
                        is_italic = bool(flags & 0b0010)

                        # Alignment heuristic
                        page_w = page.rect.width
                        cx     = (bbox[0] + bbox[2]) / 2
                        if abs(cx - page_w / 2) < page_w * 0.1:
                            alignment = "center"
                        elif cx < page_w * 0.3:
                            alignment = "left"
                        else:
                            alignment = "right"

                        tb = TextBlock(
                            text=text,
                            confidence=0.97,
                            bbox=bbox,
                            font_size=font_size,
                            font_name=font_name,
                            is_bold=is_bold,
                            is_italic=is_italic,
                            alignment=alignment,
                        )
                        page_blocks.append(tb)
                        all_blocks.append(tb)
                        page_text += text + " "

            page_data.append({
                "page_number": page_num + 1,
                "text": page_text.strip(),
                "blocks": page_blocks,
            })

        doc.close()

        raw_text = "\n\n".join(
            f"── Page {p['page_number']} ──\n{p['text']}"
            for p in page_data
            if p["text"]
        )

        return {
            "status": "success",
            "raw_text": raw_text,
            "blocks": all_blocks,
            "pages": page_data,
            "confidence": 0.97,
            "page_count": len(page_data),
            "block_count": len(all_blocks),
            "source": "pymupdf",
        }

    # ── Office fallbacks ──────────────────────────────────────────────────────

    def _office_fallback(self, file_path: str, ext: str) -> Dict:
        if ext in (".docx", ".doc"):
            return self._extract_docx(file_path)
        if ext in (".pptx", ".ppt"):
            return self._extract_pptx(file_path)
        if ext in (".xlsx", ".xls"):
            return self._extract_xlsx(file_path)
        return self._error_result("No text found and no office fallback available.")

    def _extract_docx(self, file_path: str) -> Dict:
        try:
            from docx import Document
            doc = Document(file_path)
            blocks: List[TextBlock] = []
            lines = []

            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                is_bold = any(r.bold for r in para.runs if r.bold is not None)
                is_italic = any(r.italic for r in para.runs if r.italic is not None)
                font_size = 12.0
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                        break

                align_map = {None: "left", 0: "left", 1: "center", 2: "right", 3: "justify"}
                alignment = align_map.get(para.alignment, "left")

                tb = TextBlock(
                    text=text, confidence=0.98,
                    bbox=[0, i * 20, 100, (i + 1) * 20],
                    font_size=font_size, font_name="Calibri",
                    is_bold=is_bold, is_italic=is_italic, alignment=alignment,
                )
                blocks.append(tb)
                lines.append(text)

            return {
                "status": "success",
                "raw_text": "\n".join(lines),
                "blocks": blocks,
                "confidence": 0.98,
                "block_count": len(blocks),
                "source": "docx",
            }
        except ImportError:
            return self._error_result("python-docx not installed. Run: pip install python-docx")
        except Exception as e:
            return self._error_result(f"DOCX error: {e}")

    def _extract_pptx(self, file_path: str) -> Dict:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            blocks: List[TextBlock] = []
            slides_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_lines = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in para.runs).strip()
                        if line:
                            tb = TextBlock(
                                text=line, confidence=0.97,
                                bbox=[0, 0, 0, 0],
                                font_size=18.0, font_name="Calibri",
                            )
                            blocks.append(tb)
                            slide_lines.append(line)
                if slide_lines:
                    slides_text.append(f"── Slide {slide_num} ──\n" + "\n".join(slide_lines))

            return {
                "status": "success",
                "raw_text": "\n\n".join(slides_text),
                "blocks": blocks,
                "confidence": 0.97,
                "block_count": len(blocks),
                "source": "pptx",
            }
        except ImportError:
            return self._error_result("python-pptx not installed. Run: pip install python-pptx")
        except Exception as e:
            return self._error_result(f"PPTX error: {e}")

    def _extract_xlsx(self, file_path: str) -> Dict:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            blocks: List[TextBlock] = []
            sheets_text = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        rows.append(row_str)
                        tb = TextBlock(text=row_str, confidence=0.99, bbox=[0, 0, 0, 0])
                        blocks.append(tb)
                if rows:
                    sheets_text.append(f"── Sheet: {sheet_name} ──\n" + "\n".join(rows))

            wb.close()
            return {
                "status": "success",
                "raw_text": "\n\n".join(sheets_text),
                "blocks": blocks,
                "confidence": 0.99,
                "block_count": len(blocks),
                "source": "xlsx",
            }
        except ImportError:
            return self._error_result("openpyxl not installed. Run: pip install openpyxl")
        except Exception as e:
            return self._error_result(f"XLSX error: {e}")

    # ── PaddleOCR (images only) ───────────────────────────────────────────────

    def _get_paddle(self):
        """Lazy-load PaddleOCR so PDF/DOCX paths don't pay the init cost."""
        if self._paddle_ocr is None:
            try:
                from paddleocr import PaddleOCR
                import numpy as np
                self._np = np
                # Use first language only; PaddleOCR 2.x doesn't accept a list
                lang = self.languages[0] if self.languages else "en"
                # Only 'en', 'ch', 'latin' etc are valid PaddleOCR lang codes
                valid_paddle_langs = {"en", "ch", "latin", "arabic", "cyrillic",
                                      "devanagari", "korean", "japan"}
                if lang not in valid_paddle_langs:
                    lang = "en"
                self._paddle_ocr = PaddleOCR(
                    use_angle_cls=True,
                    lang=lang,
                    show_log=False,
                )
            except ImportError:
                return None
        return self._paddle_ocr

    def _extract_image(self, file_path: str) -> Dict:
        """Extract text from an image using PaddleOCR."""
        paddle = self._get_paddle()
        if paddle is None:
            return self._error_result(
                "PaddleOCR not installed. Run: pip install paddleocr paddlepaddle"
            )

        try:
            import numpy as np
            results = paddle.ocr(file_path, cls=True)
            blocks: List[TextBlock] = []
            raw_lines = []

            if results and results[0]:
                for line in results[0]:
                    bbox, (text, conf) = line[0], line[1]
                    if not text.strip():
                        continue

                    xs = [p[0] for p in bbox]
                    ys = [p[1] for p in bbox]
                    height = max(ys) - min(ys)
                    font_size = max(8, int(height * 0.8))

                    tb = TextBlock(
                        text=text,
                        confidence=float(conf),
                        bbox=[min(xs), min(ys), max(xs), max(ys)],
                        font_size=font_size,
                        font_name="Unknown",
                    )
                    blocks.append(tb)
                    raw_lines.append(text)

            avg_conf = float(np.mean([b.confidence for b in blocks])) if blocks else 0.0

            return {
                "status": "success",
                "raw_text": "\n".join(raw_lines),
                "blocks": blocks,
                "confidence": avg_conf,
                "block_count": len(blocks),
                "source": "paddleocr",
            }

        except Exception as e:
            return self._error_result(f"PaddleOCR error: {e}")

    # ── Formatting helper (used by app) ───────────────────────────────────────

    def format_text_with_alignment(self, blocks: List[TextBlock]) -> str:
        lines = []
        for b in blocks:
            text = b.text
            if b.is_bold:
                text = f"**{text}**"
            if b.is_italic:
                text = f"_{text}_"
            if b.alignment == "center":
                text = f"{text:^80}"
            elif b.alignment == "right":
                text = f"{text:>80}"
            lines.append(text)
        return "\n".join(lines)

    # ── Internal ──────────────────────────────────────────────────────────────

    @staticmethod
    def _error_result(message: str) -> Dict:
        return {
            "status": "error",
            "raw_text": f"[{message}]",
            "blocks": [],
            "confidence": 0.0,
            "block_count": 0,
            "source": "error",
        }
